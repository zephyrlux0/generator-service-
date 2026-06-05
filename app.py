from flask import Flask, request, send_file, jsonify
from dotenv import load_dotenv
from fpdf import FPDF
import openpyxl
from pptx import Presentation
import io
import os
from datetime import datetime

load_dotenv()
app = Flask(__name__)

# ──────────────────────────────────────
# Sécurité : vérifie le token d'accès
# ──────────────────────────────────────
def check_auth():
    auth = request.headers.get('Authorization')
    expected = f"Bearer {os.getenv('API_SECRET')}"
    if auth != expected:
        return False
    return True

# ──────────────────────────────────────
# Route : PDF simple
# ──────────────────────────────────────
@app.route('/pdf', methods=['POST'])
def generate_pdf():
    if not check_auth():
        return jsonify({'error': 'Unauthorized'}), 401

    data = request.get_json()
    title = data.get('title', 'Document')
    content = data.get('content', '')
    if not content:
        return jsonify({'error': 'Missing content'}), 400

    pdf = FPDF()
    pdf.add_page()
    if title:
        pdf.set_font("Arial", 'B', 16)
        pdf.cell(200, 10, txt=title, ln=1, align='C')
        pdf.ln(10)
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, txt=content)

    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return send_file(buf, mimetype='application/pdf',
                     as_attachment=True, download_name='document.pdf')

# ──────────────────────────────────────
# Route : Excel
# ──────────────────────────────────────
@app.route('/excel', methods=['POST'])
def generate_excel():
    if not check_auth():
        return jsonify({'error': 'Unauthorized'}), 401

    data = request.get_json()
    headers = data.get('headers', [])
    rows = data.get('rows', [])
    title = data.get('title', 'Feuille1')

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title
    if headers:
        ws.append(headers)
    for row in rows:
        ws.append(row)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return send_file(buf,
                     mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                     as_attachment=True, download_name='tableau.xlsx')

# ──────────────────────────────────────
# Route : PowerPoint
# ──────────────────────────────────────
@app.route('/ppt', methods=['POST'])
def generate_ppt():
    if not check_auth():
        return jsonify({'error': 'Unauthorized'}), 401

    data = request.get_json()
    slides = data.get('slides', [])

    prs = Presentation()
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data.get('title', '')
        slide.placeholders[1].text = slide_data.get('content', '')

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return send_file(buf,
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                     as_attachment=True, download_name='presentation.pptx')

# ──────────────────────────────────────
# Route : CV
# ──────────────────────────────────────
@app.route('/cv', methods=['POST'])
def generate_cv():
    if not check_auth():
        return jsonify({'error': 'Unauthorized'}), 401

    data = request.get_json()
    name = data.get('name', '')
    email = data.get('email', '')
    phone = data.get('phone', '')
    summary = data.get('summary', '')
    experiences = data.get('experiences', [])
    education = data.get('education', [])
    skills = data.get('skills', [])

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", 'B', 20)
    pdf.cell(200, 10, txt=name, ln=1, align='C')
    pdf.set_font("Arial", size=10)
    pdf.cell(200, 6, txt=f"{email} | {phone}", ln=1, align='C')
    pdf.ln(6)
    pdf.set_font("Arial", size=10)
    pdf.multi_cell(0, 5, txt=summary)
    pdf.ln(4)

    pdf.set_font("Arial", 'B', 13)
    pdf.cell(0, 8, 'Expérience', ln=1)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(4)
    pdf.set_font("Arial", size=10)
    for exp in experiences:
        pdf.set_font("Arial", 'B', 10)
        pdf.cell(0, 6, f"{exp.get('title')} - {exp.get('company')} ({exp.get('period')})", ln=1)
        pdf.set_font("Arial", size=10)
        pdf.multi_cell(0, 5, exp.get('description', ''))
        pdf.ln(2)

    pdf.ln(4)
    pdf.set_font("Arial", 'B', 13)
    pdf.cell(0, 8, 'Formation', ln=1)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(4)
    for edu in education:
        pdf.cell(0, 6, f"{edu.get('degree')} - {edu.get('school')} ({edu.get('year')})", ln=1)
    pdf.ln(4)

    pdf.set_font("Arial", 'B', 13)
    pdf.cell(0, 8, 'Compétences', ln=1)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(4)
    pdf.set_font("Arial", size=10)
    pdf.cell(0, 6, ', '.join(skills), ln=1)

    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return send_file(buf, mimetype='application/pdf',
                     as_attachment=True, download_name='cv.pdf')

# ──────────────────────────────────────
# Route : Facture
# ──────────────────────────────────────
@app.route('/invoice', methods=['POST'])
def generate_invoice():
    if not check_auth():
        return jsonify({'error': 'Unauthorized'}), 401

    data = request.get_json()
    client_name = data.get('client_name', 'Client')
    items = data.get('items', [])
    total = sum(it.get('quantity', 0) * it.get('unit_price', 0) for it in items)

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(0, 10, 'FACTURE', ln=1, align='C')
    pdf.ln(10)
    pdf.set_font("Arial", size=10)
    pdf.cell(0, 6, f"Client : {client_name}", ln=1)
    pdf.cell(0, 6, f"Date : {datetime.now().strftime('%d/%m/%Y')}", ln=1)
    pdf.ln(10)

    pdf.set_fill_color(230, 230, 230)
    pdf.set_font("Arial", 'B', 10)
    pdf.cell(80, 7, 'Description', 1, 0, 'C', 1)
    pdf.cell(30, 7, 'Quantité', 1, 0, 'C', 1)
    pdf.cell(30, 7, 'Prix unit.', 1, 0, 'C', 1)
    pdf.cell(40, 7, 'Total', 1, 1, 'C', 1)
    pdf.set_font("Arial", size=10)
    for item in items:
        pdf.cell(80, 6, item.get('description', ''), 1)
        pdf.cell(30, 6, str(item.get('quantity', 0)), 1, 0, 'C')
        pdf.cell(30, 6, f"{item.get('unit_price', 0):.2f}", 1, 0, 'R')
        pdf.cell(40, 6, f"{item.get('quantity', 0) * item.get('unit_price', 0):.2f}", 1, 1, 'R')
    pdf.set_font("Arial", 'B', 10)
    pdf.cell(140, 7, 'TOTAL', 1, 0, 'R')
    pdf.cell(40, 7, f"{total:.2f}", 1, 1, 'R')

    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return send_file(buf, mimetype='application/pdf',
                     as_attachment=True, download_name='facture.pdf')

# ──────────────────────────────────────
# Route : Santé
# ──────────────────────────────────────
@app.route('/')
def health():
    return 'Generator Service OK'

# ──────────────────────────────────────
# Démarrage
# ──────────────────────────────────────
if __name__ == '__main__':
    port = int(os.getenv('PORT', 5000))
    app.run(host='0.0.0.0', port=port)
